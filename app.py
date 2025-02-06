import os
import json
import io
# import pymupdf as fitz
import streamlit as st
import pandas as pd
from pptx import Presentation
from PIL import Image
from langchain.prompts import PromptTemplate
from langchain_groq import ChatGroq
from langchain.schema import SystemMessage, HumanMessage
from dotenv import load_dotenv

# Initialize LangChain with Groq

load_dotenv()

groq = ChatGroq(model="llama-3.3-70b-versatile", api_key=os.getenv("API_KEY"))

# Prompt template for analysis
prompt_template = PromptTemplate(
    input_variables=["pitch_text"],
    template="""
    Analyze the following startup pitch deck text and extract structured insights. 
    **Ensure the response is in valid JSON format.** 

    Extract the following details:

    - **Name of the Startup**
    - **Problem statement**: Clearly state the problem the startup is solving in a single line.
    - **Solution**: Provide a one-line summary of how the startup addresses the problem.
    - **Elevator pitch**: Summarize the startup's value proposition in one sentence.
    - **Founded year**: Extract the year in which the startup was founded or incorporated.
    - **Location**: Mention the city and country where the startup is based or registered.
    - **Industry**: Identify the industry or sector in which the startup operates.
    - **Stage of the startup**: Determine if the startup is in the idea stage, early-stage, growth-stage, or mature stage based on available context.
    - **Website**: Provide the startup's official website URL (if available).
    - **Market size analysis (TAM, SAM, SOM)**: 
      - **TAM (Total Addressable Market)**: One-line estimate of the full market demand.
      - **SAM (Serviceable Available Market)**: The subset of TAM that the startup can realistically serve.
      - **SOM (Serviceable Obtainable Market)**: The market share the startup aims to capture.
    
    - **Revenue channels**: How does the startup generate revenue?
    - **Customer segments**: Who are the target customers?
    - **Consumer focus**: Specify whether the startup operates in B2B, B2C, or a hybrid model.
    - **Competitive advantage**: Identify the startup’s unique strengths over competitors in one line.
    - **Product innovation**: Describe any novel technology, features, or differentiators.
    - **Investment looking for**: Specify the funding amount the startup is seeking.
    - **Team profile**: List key team members along with their roles.  Ex: [Tanmay : CEO, 15yrs + experince in texh, RAvi: cto, 10+ yrs in sales]

    

    

    - **Do not assume details** if they are not explicitly mentioned. 
    - **Each response should be a single line of meaningful text, not just keywords.** 
    - If any section is missing, return **"Not mentioned"** instead of making assumptions.
    - In your response, dont have any characters that dont confine with json rules, even '/n'. Return pure JSON object
    

    **Text to analyze:**  
    ```{pitch_text}```
    """
)


def generate_report(text):
    response = groq.invoke([SystemMessage(content="Extract startup details"), HumanMessage(content=prompt_template.format(pitch_text=text))])
    return response.content

# Extract text from PPT files
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# def extract_text_from_pdf(pdf_path):
#     # Open the PDF file
#     doc = fitz.open(pdf_path)
    
#     # Extract text from each page
#     for page_num in range(len(doc)):
#         page = doc.load_page(page_num)  # Load page
#         text = page.get_text("text")   # Extract plain text
        
    
#     return doc


def process_pitchdeck(file_path):
    if file_path.endswith(".pptx"):
        extracted_text = extract_text_from_ppt(file_path)
    # elif file_path.endswith(".pptx"):
    #     extracted_text = extract_text_from_pdf(file_path)
    else:
        return "Unsupported file format"
    
    report = generate_report(extracted_text)
    return report

# Streamlit App UI
st.set_page_config(page_title="Startup Pitch Analyzer", layout="wide")
st.title("🚀 Startup Pitch Analyzer")
st.write("Upload your pitch deck (PPTX) and get detailed insights from the AI.")

uploaded_files = st.file_uploader("Upload a Pitch Deck (PPTX)", type=["pptx"], accept_multiple_files=True)

if uploaded_files:
    all_reports = []
    for uploaded_file in uploaded_files:
        file_path = f"temp_{uploaded_file.name}"
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        st.write("🔍 Extracting, please wait...")
        response = process_pitchdeck(file_path)
        # print( response)
        

        try:
            json_data = json.loads(response)
            # print("json data:", json_data)
            all_reports.append(json_data)
            st.success("✅ Analysis Complete!")
            st.json(json_data)
        except json.JSONDecodeError:
            st.error("⚠️ Error: LLM response is not valid JSON!")
    
    if all_reports:
        df = pd.DataFrame(all_reports)
        st.write("✅ Analysis Complete! Displaying Results:")
        st.write(df)
        
        excel_file = io.BytesIO()
        df.to_excel(excel_file,index=False)
        st.download_button(
            label="Download Excel Report",
            data=excel_file,
            file_name="pitch_deck_analysis.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
